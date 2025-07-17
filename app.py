import os
import pdfplumber
from pptx import Presentation
from flask import Flask, render_template, request
import google.generativeai as genai
from dotenv import load_dotenv
import threading
import webbrowser

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

# Load Gemini API key
load_dotenv()
genai.configure(api_key=os.getenv("AIzaSyA5gQ__yqKut5Wlbds7zDcUMEBo12mMI8k"))
model = genai.GenerativeModel("gemini-1.5-flash")

def extract_text_from_pdf(path):
    full_text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                full_text += text + "\n"
    return full_text

def extract_text_from_pptx(path):
    prs = Presentation(path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    return full_text

def analyze_idea(idea_text):
    prompt = f"""
შეაფასე შემდეგი ბიზნეს იდეა 7 სტრუქტურით:

ბიზნეს იდეა:
{idea_text}

1. იდეის მოკლე რეზიუმე  
2. მიზნობრივი აუდიტორია  
3. მონეტიზაციის გზები  
4. ანალოგიური პროდუქტები ან კონკურენტები (თუ იცი)  
5. იდეის სიძლიერეები და სუსტი მხარეები  
6. გრძელვადიანი მდგრადობის პროგნოზი  
7. რეკომენდაცია იდეის გაუმჯობესებისთვის
"""
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"შეცდომა: {e}"

@app.route("/", methods=["GET", "POST"])
def index():
    result = None
    if request.method == "POST":
        text_idea = request.form.get("text_idea")
        if text_idea and text_idea.strip():
            result = analyze_idea(text_idea.strip())
        elif request.files.get("file"):
            file = request.files["file"]
            if file:
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
                file.save(filepath)
                if filepath.endswith(".pdf"):
                    text = extract_text_from_pdf(filepath)
                elif filepath.endswith(".pptx"):
                    text = extract_text_from_pptx(filepath)
                else:
                    result = "❌ მხარდაჭერილია მხოლოდ .pdf და .pptx ფაილები."
                    return render_template("index.html", result=result)
                if not text.strip():
                    result = "⚠️ ფაილში ტექსტი ვერ მოიძებნა."
                else:
                    result = analyze_idea(text)
    return render_template("index.html", result=result)

if __name__ == "__main__":
    import threading
    import webbrowser

    def open_browser():
        webbrowser.open_new("http://127.0.0.1:5000")

    threading.Timer(1.25, open_browser).start()

    os.makedirs("uploads", exist_ok=True)
    app.run(debug=False)


